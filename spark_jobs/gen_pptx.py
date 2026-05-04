"""Generate OSS Pulse presentation from NYU template."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
import os

TEMPLATE = "/Users/lizhechen/Downloads/BDAD/OSS_Pulse/report/presentation/Copy of NYU Presentation (Bold) Official Template.pptx"
OUTPUT = "/Users/lizhechen/Downloads/BDAD/OSS_Pulse/report/presentation/group9_oss_pulse.pptx"

prs = Presentation(TEMPLATE)

# Remove existing slides (template has blank ones)
while len(prs.slides) > 0:
    rId = prs.slides._sldIdLst[0].rId
    prs.part.drop_rel(rId)
    prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

def add_title_body(title, body, layout_idx=5):
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    slide.placeholders[0].text = title
    slide.placeholders[1].text = body
    return slide

def add_title_only(title, layout_idx=7):
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    slide.placeholders[0].text = title
    return slide

# ── Slide 1: Title ────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[0])  # TITLE layout
slide.placeholders[0].text = "Big Data Analytics Symposium - Spring 2026"
slide.placeholders[1].text = ""
slide.placeholders[2].text = (
    "Analytics Project:  OSS Pulse: From Hype to Health\n\n"
    "Team Name:  Group 9 - OSS Pulse\n\n"
    "Team:\n"
    "Jhe Chen Li (jl17797) - Bo Yu (by2566)"
)

# ── Slide 2: Abstract + Platform ──────────────────────────────────────────────
add_title_body("OSS Pulse", (
    "Abstract:\n"
    "We analyze the health of AI open-source repositories by combining "
    "GitHub Archive event data (485+ days, ~146M events), Hugging Face Hub "
    "model metadata (2.8M models), and PyPI monthly download statistics "
    "(46 libraries). Using Apache Spark on NYU Dataproc, we compute a "
    "composite health score across three dimensions: community engagement "
    "(GitHub), model adoption (Hugging Face), and engineering usage (PyPI).\n\n"
    "We also track 5-year Q1 era comparison (2022-2026) to measure the impact "
    "of coding agents on the open-source ecosystem.\n\n\n"
    "Platform(s) where the application runs:\n"
    "Google Cloud Dataproc (Spark 3.x, YARN), HDFS"
))

# ── Slide 3: Motivation ──────────────────────────────────────────────────────
add_title_body("OSS Pulse\nMotivation", (
    "Who are the users of this analytic?\n"
    "ML engineers, researchers, project maintainers, VCs evaluating OSS frameworks\n\n"
    "Who will benefit from this analytic?\n"
    "Anyone deciding which AI framework to adopt or contribute to\n\n"
    "Why is this analytic important?\n"
    "GitHub stars alone are misleading - hype does not equal health. "
    "This project triangulates real adoption signals across three independent "
    "dimensions (GitHub activity, HF model downloads, PyPI installs) to "
    "distinguish genuinely healthy projects from viral-but-hollow ones.\n\n"
    "Era analysis: 5-year Q1 timeline (2022-2026) tracking coding agent impact "
    "on the entire OSS ecosystem."
))

# ── Slide 4: Goodness ────────────────────────────────────────────────────────
add_title_body("OSS Pulse\nGoodness", (
    "What steps were taken to assess the 'goodness' of the analytic?\n\n"
    "1. Completeness: GH Archive covers 334/365 days (bad-date filtered), "
    "supplemented with 151 additional days (Dec 2025 - Apr 2026)\n\n"
    "2. Volume: ~146M total GitHub events across 485+ days\n\n"
    "3. Multi-source validation: Health score cross-references 3 independent "
    "data sources (GH Archive, HF Hub, PyPI) - no single source can game the score\n\n"
    "4. PR double-counting fix: Used countDistinct(pr_number) instead of event "
    "counts to avoid lifecycle event inflation\n\n"
    "5. Contributor analysis validated: top1_push_ratio interpreted with "
    "pr_contributors to distinguish centralized merge from single-maintainer risk"
))

# ── Slide 5: Data Sources ────────────────────────────────────────────────────
add_title_body("OSS Pulse\nData Sources", (
    "Name:        GitHub Archive (GH Archive)\n"
    "Description: Hourly event logs - Watch/Fork/Push/PR/Issue events\n"
    "Size of data: ~146M events, 485+ days, ~99 GB cleaned Parquet\n\n"
    "Name:        Hugging Face Hub\n"
    "Description: Model metadata snapshot (April 2026) via huggingface_hub API\n"
    "Size of data: 2,815,064 models, 229 MB\n\n"
    "Name:        PyPI (BigQuery)\n"
    "Description: Monthly download counts per AI/ML library\n"
    "Size of data: 46 libraries x 12 months"
))

# ── Slide 6: Data Sample - GH Archive ────────────────────────────────────────
add_title_body("OSS Pulse\nData Sample: GitHub Archive", (
    "event_type    | event_date | actor_login | repo_name              | push_size\n"
    "WatchEvent    | 2025-03-15 | user_abc    | huggingface/transformers | null\n"
    "PushEvent     | 2025-03-15 | dev_xyz     | pytorch/pytorch          | 3\n"
    "PullRequestEvent | 2025-03-15 | contrib_1 | tensorflow/tensorflow  | null\n"
    "IssuesEvent   | 2025-03-15 | reporter_2  | microsoft/onnxruntime    | null\n"
    "ForkEvent     | 2025-03-15 | forker_3    | openai/whisper           | null\n\n"
    "Schema: event_type, event_date, actor_login, repo_name, push_size,\n"
    "push_distinct_size, commit_count, pr_merged, issue_state, pr_number,\n"
    "payload_action"
))

# ── Slide 7: Data Sample - HF Hub ────────────────────────────────────────────
add_title_body("OSS Pulse\nData Sample: Hugging Face Hub", (
    "model_id                     | library_name | pipeline_tag    | downloads  | likes\n"
    "meta-llama/Llama-3-8B        | transformers | text-generation | 1,234,567  | 8,920\n"
    "stabilityai/stable-diffusion | diffusers    | image-gen       | 890,123    | 5,430\n"
    "openai/whisper-large-v3      | transformers | auto-speech     | 456,789    | 3,210\n\n"
    "Schema: model_id, author, library_name, pipeline_tag,\n"
    "downloads, likes, created_at, parameter_count, has_safetensors"
))

# ── Slide 8: Data Sample - PyPI ──────────────────────────────────────────────
add_title_body("OSS Pulse\nData Sample: PyPI", (
    "project      | month      | downloads\n"
    "transformers | 2025-01-01 | 12,345,678\n"
    "torch        | 2025-01-01 | 45,123,456\n"
    "diffusers    | 2025-02-01 | 3,456,789\n"
    "accelerate   | 2025-03-01 | 2,134,500\n"
    "datasets     | 2025-04-01 | 4,890,123\n\n"
    "Schema: project, month (YYYY-MM-01), downloads"
))

# ── Slide 9: Design Diagram ──────────────────────────────────────────────────
add_title_body("OSS Pulse\nDesign Diagram", (
    "Data Ingestion              Processing (Spark)           Output\n"
    "----------------            ------------------           ------\n"
    "GH Archive (HDFS) --+\n"
    "                    +---> Job 01: AI repo daily metrics\n"
    "HF Hub API ---------+     Job 02: HF + GH join          health_score.csv\n"
    "                    +---> Job 03: Three-way health score\n"
    "PyPI BigQuery ------+     Job 04: Top 1000 repos         Tableau vizs\n"
    "                         Job 05: AI vs General\n"
    "                         Job 06: Star hype detection\n"
    "                         Job 07: Contributor health\n"
    "                         Job 08: Era comparison (5yr Q1)\n"
    "                         Job 09: Repo deep dive\n\n"
    "Tools: Python (ingest) | PySpark (analytics) | Scala Spark (HF) | Tableau\n\n"
    "Health Score = sum of log1p(metric) x weight across 6 signals\n"
    "Extends OpenSSF Criticality Score with HF + PyPI adoption dimensions"
))

# ── Slide 10: Code Challenge 1 (jl17797) ─────────────────────────────────────
add_title_body("OSS Pulse\nCode Challenge: jl17797 - HDFS Quota & Rolling Pipeline", (
    "Challenge 1 - Data freshness gap:\n"
    "HF Hub snapshot is April 2026, but GH Archive only covered through "
    "Nov 2025. Needed to ingest 3,624 hourly files.\n\n"
    "Solution: Stream directly to HDFS via curl | hdfs dfs -put -f -\n"
    "One-time hdfs dfs -ls + local grep for resume-safe skip.\n\n"
    "Challenge 2 - HDFS quota: 500 GB hard limit\n"
    "5 years x Q1 raw = ~700 GB needed, exceeds quota.\n\n"
    "Solution - Rolling pipeline orchestration:\n"
    "Process one year at a time: download -> Spark clean -> delete raw -> next year.\n"
    "Cleaned Parquet is 8x smaller than raw gz.\n\n"
    "Outcome: 5-year Q1 timeline (2022-2026) ingested within 500 GB quota."
))

# ── Slide 11: Code Challenge 2 (jl17797) ─────────────────────────────────────
add_title_body("OSS Pulse\nCode Challenge: jl17797 - Era Comparison at Scale", (
    "Challenge 3 - 90 GB across 5 eras, 48 GB executor memory:\n"
    "Naive approach (cache all -> 4 aggregations) failed: cache spills to disk,\n"
    "YARN kills job after 1.5h.\n\n"
    "Solution - Per-era independent processing:\n"
    "Process each era (15-23 GB) independently, cache fits in memory,\n"
    "merge 5 small result sets at the end.\n\n"
    "                    | Union-first      | Per-era\n"
    "Peak memory   | 90 GB (exceeds)  | 15-23 GB (fits)\n"
    "Failure blast   | All restart        | Only failed era\n"
    "Runtime          | 2-3h (killed)    | 40-50 min\n\n"
    "Challenge 4 - PR double counting:\n"
    "Used countDistinct(pr_number) instead of event count.\n"
    "Filter payload_action='closed' AND pr_merged=True for merge rate."
))

# ── Slide 12: Code Challenge 3 (by2566) ──────────────────────────────────────
add_title_body("OSS Pulse\nCode Challenge: by2566 - Billion-Record Shuffle", (
    "Challenge: GH Archive partitioned by event_date (334 partitions),\n"
    "but analysis requires aggregation by repo_name.\n"
    "Computing top-1000 repos across 1B+ events forces full cross-partition shuffle.\n\n"
    "Scale: 72 GB input -> 83 GB shuffle -> 1B+ records, 3+ hours\n\n"
    "Solution: Tuned spark.sql.shuffle.partitions (400 -> 200).\n"
    "Structured downstream jobs to read from pre-aggregated Job 04 output\n"
    "instead of re-shuffling raw events.\n\n"
    "Additional optimizations:\n"
    "- .cache() on computed DataFrames before multiple write actions\n"
    "- Replaced Python UDFs with native F.col().isin()\n"
    "- Column pruning: select only needed columns for each aggregation"
))

# ── Slide 13: Results ────────────────────────────────────────────────────────
add_title_body("OSS Pulse\nResults", (
    "1. Hype != Health: Ollama has 44,259 stars but only 241 HF downloads.\n"
    "   DeepSeek-R1/V3 rank #2/#3 in stars but have no sustainable tool ecosystem.\n\n"
    "2. Quiet Powerhouses exist: sentence-transformers has 327,014 HF downloads\n"
    "   per star - 3x more efficient than transformers. Stars undercount real adoption.\n\n"
    "3. Bus Factor risk is real: hiyouga/llama-factory has top1_push_ratio 0.961\n"
    "   (but 123 PR contributors = centralized merge, not single maintainer).\n"
    "   geekan/metagpt: 0.929 with only 22 PR contributors = genuine risk.\n\n"
    "Key insight: GitHub stars are the worst health indicator.\n"
    "PyPI + HF downloads together better predict sustained ecosystem health.\n"
    "Push contributor concentration reveals single-maintainer risk invisible to stars."
))

# ── Slide 14: Obstacles ──────────────────────────────────────────────────────
add_title_body("OSS Pulse\nObstacles", (
    "1. HDFS permission boundary\n"
    "   jl17797 has read-only access to by2566's directory; cannot write there.\n"
    "   Workaround: all outputs go to jl17797's HDFS; Spark jobs union both paths.\n\n"
    "2. Data freshness mismatch\n"
    "   GH Archive: 2025 only; HF Hub: April 2026 snapshot.\n"
    "   Required supplemental ingestion of 151 days (Dec 2025 - Apr 2026),\n"
    "   streaming ~3,600 files to HDFS. ~580 corrupt files in April detected\n"
    "   and handled via ignoreCorruptFiles.\n\n"
    "3. YARN resource contention on shared cluster\n"
    "   Other users' jobs caused executor preemption and OOM kills.\n"
    "   Solution: sequential execution, conservative resource configs,\n"
    "   per-era job design to minimize runtime and memory footprint."
))

# ── Slide 15: Summary ────────────────────────────────────────────────────────
add_title_body("OSS Pulse\nSummary", (
    "- Built end-to-end Spark pipeline combining 3 data sources at ~146M-event scale\n\n"
    "- Computed composite health score for 36 AI repos across community,\n"
    "  adoption, and engineering dimensions\n\n"
    "- Analyzed contributor health and bus-factor risk for top 1,000 repos\n\n"
    "- 5-year era comparison (2022-2026 Q1) tracking coding-agent impact\n\n"
    "- Deep-dive: ~250 repos (AI vs Non-AI) across 5 eras\n\n"
    "- Key takeaway: GitHub stars are a noisy signal;\n"
    "  PyPI + HF downloads together better predict sustained ecosystem health.\n"
    "  Push contributor concentration reveals single-maintainer risk\n"
    "  invisible to star counts.\n\n"
    "- Framework-agnostic methodology - extensible to any open-source domain"
))

prs.save(OUTPUT)
print(f"Saved to {OUTPUT}")
print(f"Total slides: {len(prs.slides)}")
